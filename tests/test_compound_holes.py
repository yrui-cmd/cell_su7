"""A white compound must remain one shape with its hole contours intact."""
import json
from pathlib import Path
import subprocess
import sys
import tempfile
from pptx import Presentation

ROOT=Path(__file__).resolve().parents[1]
SCRIPTS=ROOT/'plugins/cell_su7/skills/cell_su7/scripts'
with tempfile.TemporaryDirectory() as directory:
    work=Path(directory)
    source=work/'holes.svg'
    source.write_text('<svg xmlns="http://www.w3.org/2000/svg" viewBox="0 0 100 100"><rect id="blue" width="100" height="100" fill="blue"/><path id="white-ring" fill="white" d="M0 0H100V100H0Z M20 20V80H80V20Z"/></svg>')
    def run(name,*args):
        subprocess.run([sys.executable,str(SCRIPTS/name),*map(str,args)],check=True,capture_output=True)
    run('prepare_geometry_cache.py','--input',source,'--output-dir',work/'cache','--job-id','holes')
    run('cull_hidden_geometry.py','--cache',work/'cache/geometry-cache.json','--state',work/'cache/drawing-state.json')
    cache=json.loads((work/'cache/geometry-cache.json').read_text(encoding='utf-8'))
    assert len(cache['atoms'])==2 and len(cache['atoms'][1]['subpaths'])==2
    output=work/'holes.pptx'
    run('run_cell_ppt_ooxml.py','--geometry-cache',work/'cache/geometry-cache.json','--output-pptx',output)
    prs=Presentation(output)
    assert len(prs.slides[0].shapes)==2
    ring=prs.slides[0].shapes[1]
    assert len(ring._element.xpath('.//a:path'))==1
    assert len(ring._element.xpath('.//a:moveTo'))==2
    assert len(ring._element.xpath('.//a:close'))==2
print('COMPOUND_HOLES_OK|contours=preserved|white-ring=one-native-shape|no-subpath-split')
