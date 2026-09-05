"""Restore original lettering through both backends without any paid API calls."""
import importlib.util
import json
from pathlib import Path
import shutil
import subprocess
import sys
import tempfile
from unittest.mock import patch
import xml.etree.ElementTree as ET
from pptx import Presentation

ROOT = Path(__file__).resolve().parents[1]
SCRIPTS = ROOT / 'plugins/cell_su7/skills/cell_su7/scripts'
spec = importlib.util.spec_from_file_location('entry', SCRIPTS / 'run_from_image.py')
entry = importlib.util.module_from_spec(spec)
spec.loader.exec_module(entry)
real_run = entry.run
fixture = ROOT / 'tests/fixtures/visibility.svg'
manifest = ROOT / 'tests/fixtures/text-manifest.json'

def check_master(raw, master):
    before = ET.parse(raw).getroot()
    after = ET.parse(master).getroot()
    label = next(node for node in after.iter() if node.get('id') == 'fixture-label')
    assert label.tag.endswith('}text') and label.text == 'Cell_ppt'
    assert label.get('font-weight') == 'bold'
    assert float(label.get('x')) == 66 and float(label.get('y')) == 50
    for node in before:
        restored = next(child for child in after if child.get('id') == node.get('id'))
        assert ET.tostring(node) == ET.tostring(restored), 'Non-text geometry changed'

with tempfile.TemporaryDirectory(prefix='cell-gd-text-') as directory:
    work = Path(directory)
    image = work / 'cleaned.png'
    image.write_bytes(b'mocked cleaned image')
    uploads = []
    def fake_run(*args):
        if Path(args[0]).name == 'vectorize_xiaomiao.py':
            uploads.append(Path(args[args.index('--input-image') + 1]))
            shutil.copyfile(fixture, args[args.index('--output-svg') + 1])
        else:
            real_run(*args)
    with patch.object(sys, 'argv', ['entry', '--input-image', str(image), '--text-manifest', str(manifest), '--output-root', str(work / 'python')]), patch.object(entry, 'run', fake_run):
        assert entry.main() == 0
    assert uploads == [image.resolve()]
    python_job = work / 'python/shibielujing1'
    check_master(python_job / 'shibielujing1-vector.svg', python_job / 'shibielujing1.svg')
    prs = Presentation(python_job / 'shibielujing1.pptx')
    assert any(getattr(shape, 'text', '') == 'Cell_ppt' for shape in prs.slides[0].shapes)

    if sys.platform == 'win32':
        scripts = work / 'scripts'
        shutil.copytree(SCRIPTS, scripts)
        shutil.copyfile(fixture, scripts / 'fixture.svg')
        (scripts / 'vectorize-xiaomiao.ps1').write_text('''param($InputImage, $OutputSvg, $EstimatedCredits, [switch]$ApproveHighCost)
Copy-Item -LiteralPath (Join-Path $PSScriptRoot 'fixture.svg') -Destination $OutputSvg
''', encoding='utf-8-sig')
        # Use the real public selector for both apps, replacing only desktop host invocation.
        (scripts / 'run_cell_ppt.ps1').write_text('''param($GeometryCache, $OutputPptx, $HostApplication, $StepDelayMs, [switch]$UseActivePresentation, [switch]$Foreground)
py -3 -X utf8 (Join-Path $PSScriptRoot 'run_cell_ppt_ooxml.py') --geometry-cache $GeometryCache --output-pptx $OutputPptx
if ($LASTEXITCODE -ne 0) { throw 'OOXML failed' }
''', encoding='utf-8-sig')
        # The Illustrator cache is the exact live-text input for its desktop runtime.
        (scripts / 'run_cell_lct.ps1').write_text('''param($InputSvg, $WorkDir, $OutputAi, $OutputPng, $Placement, $MaxWidthFraction, $MaxHeightFraction, $DelayMs, $MinBatchSize, $MaxBatchSize)
py -3 -X utf8 (Join-Path $PSScriptRoot 'prepare_illustrator_cache.py') --input $InputSvg --output-dir $WorkDir --job-id text-test
if ($LASTEXITCODE -ne 0) { throw 'Cache failed' }
''', encoding='utf-8-sig')
        for app in ['ppt', 'ai']:
            proc = subprocess.run(['powershell', '-NoProfile', '-ExecutionPolicy', 'Bypass', '-File', str(scripts / 'run_cell_su7.ps1'), '-InputImage', str(image), '-TextManifest', str(manifest), '-OutputRoot', str(work / app), '-Application', app], capture_output=True, text=True)
            assert proc.returncode == 0, proc.stderr
            job = work / app / 'shibielujing1'
            check_master(job / 'shibielujing1-vector.svg', job / 'shibielujing1.svg')
            if app == 'ppt':
                prs = Presentation(job / 'shibielujing1.pptx')
                assert any(getattr(shape, 'text', '') == 'Cell_ppt' for shape in prs.slides[0].shapes)
            else:
                cache = json.loads((job / '.cell-lct-internal/live-cache/geometry-cache.json').read_text(encoding='utf-8'))
                assert any(atom.get('kind') == 'text' for atom in cache['atoms'])
print('TEXT_RESTORATION_OK|python=pptx-live-text|windows=ppt,ai-cache|geometry=preserved|api=mocked|desktop=not-tested')
