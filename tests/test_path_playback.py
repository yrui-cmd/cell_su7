from pathlib import Path
from tempfile import TemporaryDirectory
from zipfile import ZipFile
import sys
from lxml import etree as E
from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE
sys.path.insert(0, str(Path(__file__).resolve().parents[1] / 'plugins/cell_su7/skills/cell_su7/scripts'))
from prepare_path_playback import prepare, NS

with TemporaryDirectory() as tmp:
    root = Path(tmp)
    p = Presentation()
    s = p.slides.add_slide(p.slide_layouts[6])
    s.shapes.add_shape(MSO_SHAPE.DONUT, 0, 0, 1000000, 1000000)
    s.shapes.add_textbox(100, 100, 1000000, 1000000).text = 'Text position'
    p.save(root/'source.pptx')
    assert prepare(root/'source.pptx', root/'hidden.pptx') == 2
    with ZipFile(root/'source.pptx') as a, ZipFile(root/'hidden.pptx') as b:
        for name in a.namelist():
            if name != 'ppt/slides/slide1.xml':
                assert a.read(name) == b.read(name)
            else:
                before, after = E.fromstring(a.read(name)), E.fromstring(b.read(name))
                for shape in after.findall('p:cSld/p:spTree/p:sp', NS):
                    assert shape.find('p:nvSpPr/p:cNvPr', NS).attrib.pop('hidden') == '1'
                assert E.tostring(before) == E.tostring(after)
    try:
        prepare(root/'source.pptx', root/'source.pptx')
        raise AssertionError('overwrite allowed')
    except ValueError:
        pass
    try:
        prepare(root/'hidden.pptx', root/'bad.pptx')
        raise AssertionError('hidden helper objects accepted')
    except ValueError:
        pass
print('PATH_PLAYBACK_OK|geometry-colors-text=unchanged|overwrite=blocked|hidden-source=rejected')
