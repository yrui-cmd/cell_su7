#!/usr/bin/env python3
"""Deterministic tests for the shared cell_su7 core and OOXML backend."""

from __future__ import annotations

import json
import subprocess
import sys
import tempfile
import zipfile
from pathlib import Path

from pptx import Presentation
from pptx.util import Inches


ROOT = Path(__file__).resolve().parents[1]
SCRIPTS = ROOT / "plugins" / "cell_su7" / "skills" / "cell_su7" / "scripts"


def run(*args: object, check: bool = True) -> subprocess.CompletedProcess[str]:
    return subprocess.run([str(value) for value in args], text=True, capture_output=True, check=check)


def main() -> int:
    with tempfile.TemporaryDirectory(prefix="cell-ppt-test-") as raw:
        temp = Path(raw)
        cache = temp / "cache"
        run(sys.executable, SCRIPTS / "prepare_geometry_cache.py", "--input", ROOT / "tests" / "fixtures" / "editable.svg", "--output-dir", cache, "--job-id", "shibielujing1")
        run(sys.executable, SCRIPTS / "cull_hidden_geometry.py", "--cache", cache / "geometry-cache.json", "--state", cache / "drawing-state.json")
        payload = json.loads((cache / "geometry-cache.json").read_text(encoding="utf-8"))
        if payload["schema_version"] != 3 or not any(atom["kind"] == "text" for atom in payload["atoms"]):
            raise AssertionError("Shared geometry cache contract failed")

        source = temp / "source.pptx"
        prs = Presentation()
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        sentinel = slide.shapes.add_textbox(Inches(0.1), Inches(0.1), Inches(1), Inches(0.3))
        sentinel.name = "PREEXISTING_SENTINEL"
        sentinel.text = "keep"
        prs.save(source)

        output = temp / "output.pptx"
        proc = run(sys.executable, SCRIPTS / "run_cell_ppt_ooxml.py", "--geometry-cache", cache / "geometry-cache.json", "--input-pptx", source, "--output-pptx", output, "--slide-index", 1)
        result = json.loads(proc.stdout.strip().splitlines()[-1])
        if not result["ok"] or result["native_object_count"] < 4:
            raise AssertionError("OOXML native-object creation failed")
        reopened = Presentation(output)
        names = [shape.name for shape in reopened.slides[0].shapes]
        editable_text = [shape.text for shape in reopened.slides[0].shapes if getattr(shape, "has_text_frame", False)]
        if "PREEXISTING_SENTINEL" not in names or "cell_su7" not in editable_text:
            raise AssertionError("Existing-object or editable-text preservation failed")
        with zipfile.ZipFile(output) as package:
            if any(name.startswith("ppt/media/") for name in package.namelist()):
                raise AssertionError("Raster media was inserted into editable output")

        install_root = temp / "skills"
        run(sys.executable, ROOT / "install.py", "--destination", install_root)
        if not (install_root / "cell_su7" / "SKILL.md").is_file():
            raise AssertionError("Cross-platform installer failed")

        gate = run(
            sys.executable,
            SCRIPTS / "vectorize_xiaomiao.py",
            "--input-image", temp / "not-uploaded.png",
            "--output-svg", temp / "never-created.svg",
            "--estimated-credits", 2,
            check=False,
        )
        if gate.returncode == 0 or "CREDIT_CONFIRM_REQUIRED" not in (gate.stdout + gate.stderr):
            raise AssertionError("Pre-upload credit confirmation gate failed")

    print("CROSS_PLATFORM_OK|core=shared|ooxml=editable|raster=absent|existing=preserved|credit_gate=preupload")
    return 0


if __name__ == "__main__":
    raise SystemExit(main())
