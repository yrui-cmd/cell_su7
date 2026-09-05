"""Coordinate, anchor and rotation regressions independent of image generation."""
import importlib.util
import math
from pathlib import Path
import xml.etree.ElementTree as ET
from pptx import Presentation
from pptx.enum.text import PP_ALIGN

SCRIPTS = Path(__file__).resolve().parents[1] / 'plugins/cell_su7/skills/cell_su7/scripts'
def load(name):
    spec = importlib.util.spec_from_file_location(name, SCRIPTS / f'{name}.py')
    module = importlib.util.module_from_spec(spec)
    spec.loader.exec_module(module)
    return module

merge = load('merge_live_text')
render = load('run_cell_ppt_ooxml')
cache = load('prepare_geometry_cache')
item = {'id': 'rotated', 'content': 'Top\nBottom', 'x': 300, 'y': 320, 'coordinate_space': 'pixels', 'font_size': 20, 'rotation': 90, 'text_anchor': 'middle'}
mapped = merge.source_coordinates(item, {'source_canvas': {'width': 1200, 'height': 800}})
assert mapped['x'] == .25 and mapped['y'] == .4 and mapped['font_size'] == .025
root = ET.fromstring('<svg xmlns="http://www.w3.org/2000/svg" viewBox="10 20 600 400"/>')
root.append(merge.add_text(root, mapped, (10, 20, 600, 400)))
atoms = cache.collect_atoms(root)
assert len(atoms) == 2
assert atoms[0]['text']['position'] == [160, 180]
assert atoms[1]['text']['position'] == [148, 180], atoms[1]['text']['position']
assert atoms[0]['text']['rotationDegrees'] == 90

prs = Presentation()
slide = prs.slides.add_slide(prs.slide_layouts[6])
for anchor, alignment in [('start', PP_ALIGN.LEFT), ('middle', PP_ALIGN.CENTER), ('end', PP_ALIGN.RIGHT)]:
    atom = {'objectName': anchor, 'text': {'position': [160, 180], 'fontSize': 10, 'contents': 'Label', 'textAnchor': anchor, 'rotationDegrees': 90}}
    render.add_text(slide, atom, (2, 30, 40, 10, 20))
    shape = slide.shapes[-1]
    assert shape.text_frame.paragraphs[0].alignment == alignment
    assert shape.text_frame.word_wrap is False
    # Reconstruct the baseline anchor after Office's centre-based rotation.
    width = shape.width / render.EMU_PER_PT
    height = shape.height / render.EMU_PER_PT
    local_x = {'start': 0, 'middle': width / 2, 'end': width}[anchor]
    local_y = 20 * 1.05
    cx = shape.left / render.EMU_PER_PT + width / 2
    cy = shape.top / render.EMU_PER_PT + height / 2
    x = cx - (local_y - height / 2)
    y = cy + (local_x - width / 2)
    assert math.isclose(x, 330, abs_tol=.001) and math.isclose(y, 360, abs_tol=.001), (anchor, x, y)
print('TEXT_POSITIONS_OK|pixel-mapping=exact|viewbox-origin=preserved|multiline-rotation=shared-pivot|anchors=start,middle,end')
